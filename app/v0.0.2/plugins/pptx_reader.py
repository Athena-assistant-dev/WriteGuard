from pptx import Presentation
from plugins.base_plugin import WritePlugin

class PptxReaderPlugin(WritePlugin):
    extensions = [".pptx"]

    def read(self, filepath):
        try:
            prs = Presentation(filepath)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)
        except Exception as e:
            return f"[PPTX READ ERROR] {e}"

    def post_write(self, filepath, old_content=None, new_content=None):
        print(f"[PPTX PLUGIN] Wrote: {filepath}")

# For standalone execution/testing
if __name__ == "__main__":
    plugin = PptxReaderPlugin()